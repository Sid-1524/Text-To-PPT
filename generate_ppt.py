import os
import re
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configuration
GITHUB_TOKEN = os.getenv("GITHUB_TOKEN")
ENDPOINT = "https://models.github.ai/inference"
MODEL_NAME = "openai/gpt-4.1"

MAX_SLIDES = 7
MAX_BULLET_CHARS = 180
MAX_TOTAL_CHARS = 1500
FONT_SIZE_TITLE = Pt(40)
FONT_SIZE_CONTENT = Pt(20)

def generate_slides_content(topic):
    """Generate presentation content using AI model"""
    client = OpenAI(base_url=ENDPOINT, api_key=GITHUB_TOKEN)
    
    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=[
            {"role": "system", "content": """Create detailed technical content with:
             - Full sentence explanations
             - Concrete examples/data
             - 5 substantial points per slide"""},
            {"role": "user", "content": f"""Create a 7-slide structure about {topic} with:
                ## [Section Title]
                - Detailed point with technical specifications
                - Example: 'CNNs use 3x3 kernels with ReLU activation'
                - Continue for 5 substantial points per slide"""}
        ],
        temperature=0.7,
        top_p=0.9,
        max_tokens=2000
    )
    return response.choices[0].message.content

def parse_content(content):
    """Parse AI response into structured slides"""
    slides = []
    current_slide = {"title": "", "points": []}
    
    for line in content.split('\n'):
        line = line.strip()
        if line.startswith('## '):
            # Clean title from numbering
            title = re.sub(r'^(Slide\s*\d+:|[\d.]+)\s*', '', line[3:], flags=re.IGNORECASE).strip()
            if current_slide["title"]:
                slides.append(current_slide)
                current_slide = {"title": "", "points": []}
            current_slide["title"] = title
        elif line.startswith('- '):
            current_slide["points"].append(line[2:].strip())
    
    if current_slide["title"]:
        slides.append(current_slide)
    
    return slides[:MAX_SLIDES]

def optimize_content(points):
    """Ensure content fits slide canvas"""
    optimized = []
    total_chars = 0
    
    for point in points:
        # Truncate long points but preserve sentence endings
        if len(point) > MAX_BULLET_CHARS:
            point = point[:MAX_BULLET_CHARS].rsplit('.', 1)[0] + '...'
        
        if total_chars + len(point) <= MAX_TOTAL_CHARS:
            optimized.append(point)
            total_chars += len(point)
        else:
            remaining = MAX_TOTAL_CHARS - total_chars
            if remaining > 40:  # Minimum meaningful content
                optimized.append(point[:remaining-3] + '...')
            break
    
    return optimized

def create_presentation(topic, slides):
    """Generate PowerPoint file"""
    prs = Presentation()
    
    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    title_shape.text = topic
    title_shape.text_frame.paragraphs[0].font.size = FONT_SIZE_TITLE
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Content Slides
    for slide in slides:
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Slide Title
        title_box = content_slide.shapes.title
        title_box.text = slide["title"]
        title_box.text_frame.paragraphs[0].font.size = FONT_SIZE_TITLE
        title_box.text_frame.paragraphs[0].font.bold = True
        
        # Content Body
        content_frame = content_slide.placeholders[1].text_frame
        content_frame.clear()
        
        for point in optimize_content(slide["points"]):
            p = content_frame.add_paragraph()
            p.text = point
            p.font.size = FONT_SIZE_CONTENT
            p.level = 0
            p.space_after = Pt(6)  # Tight spacing

    filename = f"{topic.replace(' ', '_')}_presentation.pptx"
    prs.save(filename)
    return filename

def main():
    """Main execution flow"""
    topic = input("Enter presentation topic: ").strip()
    
    if not GITHUB_TOKEN:
        print("Error: GITHUB_TOKEN not found in environment")
        return
    
    try:
        content = generate_slides_content(topic)
        slides = parse_content(content)
        output_file = create_presentation(topic, slides)
        print(f"Successfully created: {output_file}")
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")

if __name__ == "__main__":
    main()
